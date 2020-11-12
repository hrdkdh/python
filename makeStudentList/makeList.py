import os
import pandas as pd
from time import sleep
from pptx.util import Pt
from pptx import Presentation

filesave_root_path = None
filesave_path = None

def checkClipboard():
    print("===================================================================================================")     
    print("엑셀 양식에 교육생 정보 데이터를 입력한 다음 클립보드에 복사해 주세요.")
    print("[명단작성 데이터 양식.xlsx] 파일을 참고하세요")
    _ = input("복사가 완료되면 엔터키를 눌러 주세요.")
    
    try:
        df = pd.read_clipboard()
        if df is None or len(df) < 1 or "성명" not in df.columns:
            print("교육생 정보가 클립보드로 복사되지 않았습니다.")
            sleep(2)
            checkClipboard()
    except:
        print("데이터가 복사되지 않았습니다.")
        sleep(2)
        checkClipboard()
    return df
    
def makePPT(pic_image_resized_path, cha_name):
    if cha_name == "" or cha_name == None:
        print("차수명을 입력하지 않았습니다.")
        cha_name = input("차수명을 입력해 주십시오.")
        makePPT(pic_image_resized_path, cha_name)

    print("차수명 : " + cha_name)
    df = checkClipboard()
    print("교육생 명단을 PPT로 작성하는 중...")

    if len(df) <= 30:
        post_file_name = "30"
    elif len(df) <= 36:
        post_file_name = "36"
    elif len(df) <= 42:
        post_file_name = "42"
    prs = Presentation("./ppt_master/master_"+post_file_name+".pptx")
    slide = prs.slides[0]
    for student in df.iloc:
        for shape in slide.shapes:
            if shape.has_text_frame:
                this_paragraph = shape.text_frame.paragraphs[0]
                if str(student["조"]).strip() == this_paragraph.text.strip()[2:3] and str(student["출력순서"]).strip() == this_paragraph.text.strip()[4:5]:
                    this_label = this_paragraph.text.strip()[0:2]
                    if this_label == "사진":
                        files = os.listdir("./"+pic_image_resized_path)
                        for f in files:
                            if student["휴대폰"] in f:
                                slide.shapes.add_picture("./"+pic_image_resized_path+f, shape.left, shape.top, shape.width, shape.height)
                                break
                        this_paragraph.text = ""

            elif shape.has_table:
                for i in range(0, 7):
                    cell = shape.table.rows[i].cells[0]
                    this_table_paragraph = cell.text_frame.paragraphs[0]
                    this_label = this_table_paragraph.text.strip()[0:2]
                    if str(student["조"]).strip() == this_table_paragraph.text.strip()[2:3] and str(student["출력순서"]).strip() == this_table_paragraph.text.strip()[4:5]:
                        if this_label == "이름":
                            this_table_paragraph.text = str(student["성명"])
                            this_table_paragraph.font.bold = True
                        elif this_label == "나이":
                            this_table_paragraph.text = str(student["나이"])
                        elif this_label == "대학":
                            this_table_paragraph.text = str(student["대학명"])
                        elif this_label == "전공":
                            this_table_paragraph.text = str(student["학부전공"])
                        elif this_label == "지역":
                            this_table_paragraph.text = str(student["거주지_시"])
                        elif this_label == "전화":
                            this_table_paragraph.text = str(student["휴대폰"])
                        elif this_label == "숙소":
                            this_table_paragraph.text = str(student["숙소"])
                        this_table_paragraph.font.size = Pt(8)

    makeDownloadDirectory()
    prs.save(filesave_path+"교육생 명부_"+cha_name+".pptx")
    print("교육생 명단 작성완료")
    print("다음 폴더에 PPT명단 파일을 저장하였습니다 : " + filesave_path+"교육생 명부_"+cha_name+".pptx")

def makeDownloadDirectory():
    global filesave_root_path, filesave_path

    #폴더 생성 및 PPT 생성을 위한 정보
    filesave_root_path = "results/"
    filesave_path = filesave_root_path+"명단작성 결과/"

    for dir_path in [filesave_root_path, filesave_path]:
        try:
            if not(os.path.isdir("./"+dir_path)):
                os.makedirs(os.path.join("./"+dir_path))
        except OSError as e:
            if e.errno != errno.EEXIST:
                print(dir_path + " : 폴더 생성에 실패하였습니다.")
                raise