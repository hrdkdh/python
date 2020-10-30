import os
import sys
import cv2
import urllib
import requests
import numpy as np
import pandas as pd
from shutil import copyfile
from pptx import Presentation
from pptx.util import Pt
from datetime import datetime
from bs4 import BeautifulSoup as bs

login_id = "edu" # 관리자 id
login_pw = "12345!" # 관리자 password 
cha_name = "20-005차 포스코 취업아카데미(포항5차)" #다운받고자 하는 차수명(정확해야 함)
image_path = "images/"
download_path = image_path+"{:%Y%m%d%H%M%S}/".format(datetime.now())
resized_path = download_path+"resized/"
image_resize_size = [800, 600] #height, width / 이미지 비율은 4:3으로 고정

def checkInit():
    check = True
    df = pd.read_clipboard()
    if len(df) < 30:
        printInitError()
        check = False
    return check

def printInitError():
    print("교육생 정보가 클립보드로 복사되지 않아 실행을 중단합니다.")
    print("입과자 명부 엑셀 파일에서 데이터를 클립보드로 복사해 주세요.")
    print("조, 출력순서, 성명, 휴대폰, 나이, 대학명, 학부전공, 졸업, 거주지, 숙소 정보가 복사되어야 합니다.")

def makeDownloadDirectory(download_path):
    try:
        if not(os.path.isdir("./"+download_path)):
            os.makedirs(os.path.join("./"+download_path))
    except OSError as e:
        if e.errno != errno.EEXIST:
            print("이미지 다운로드 폴더 생성에 실패하였습니다.")
            raise

def faceRecognition(img):
    model = "res10_300x300_ssd_iter_140000.caffemodel"
    config = "deploy.prototxt"
    net = cv2.dnn.readNet(model, config)

    if net.empty():
        print("얼굴인식에 필요한 소스파일이 없습니다.")
    else:
        blob = cv2.dnn.blobFromImage(img, 1, (300, 300), (104, 177, 123))
        net.setInput(blob)
        out = net.forward()

        detect = out[0, 0, :, :]
        (h, w) = img.shape[:2]

        #200개 중 가장 확률 높은 하나만 출력 (i=0)
        # for i in range(detect.shape[0]):
        confidence = detect[0, 2]
        x1 = int(detect[0, 3] * w)
        y1 = int(detect[0, 4] * h)
        x2 = int(detect[0, 5] * w)
        y2 = int(detect[0, 6] * h)

        cv2.rectangle(img, (x1, y1), (x2, y2), (0, 255, 0))

    return img

def downloadStudentImages(login_id, login_pw, cha_name):
    print("youth.posco.com에서 이미지를 다운받는 중...")
    base_url = "http://youth.posco.com/posco/_owner/"
    login_url = base_url+"index.php?act=login"
    login_data = {
        "wd_id": login_id,
        "wd_pw": login_pw
    }
    with requests.Session() as s:
        login_req = s.post(login_url, data=login_data)
        if login_req.status_code != 200:
            print(login_req.status_code)
            print("관리자 화면 로그인에 실패하였습니다.")
            sys.exit()
        cha_list_data = s.get(base_url+"index.php?mod=lecture&act=main&cate=&sField=&sValue="+cha_name)
        soup = bs(cha_list_data.text, "html.parser")
        table = soup.select("table")
        strongs = table[1].select("strong")
        for strong in strongs:
            if "격 : " in strong.get_text() and "합" in strong.get_text()[:1]:
                href = strong.parent.attrs["href"]
                student_page_url=href[2:len(href)]
                break
        student_list_data = s.get(base_url+student_page_url)
        image_soup = bs(student_list_data.text, "html.parser")
        image_links = image_soup.select("a")
        for href in image_links:
            if href.get_text() is not None and "증명사진" in href.get_text():
                this_href = href.attrs["href"]
                this_image = s.get(base_url+this_href[2:len(this_href)])
                this_ext = this_image.headers["Content-Disposition"].split(".")[1]
                this_image_name = href.parent.parent.findAll("td")[3].find("a").text.replace("/", "_")+"_"+href.parent.parent.findAll("td")[4].find("strong").text
                open("./"+download_path+this_image_name+"."+this_ext, "wb").write(this_image.content)

def cropImages(download_path):
    print("이미지를 4:3 비율로 자르고 다듬는 중...")
    makeDownloadDirectory(resized_path)
    try:
        os.remove("./"+download_path+"temp")
    except:
        pass
    files = os.listdir("./"+download_path)
    for f in files:
        if len(f.split("."))<2: #폴더일 경우는 처리하지 않음
            continue

        dst = "./"+download_path+"temp"
        copyfile("./"+download_path+f, dst)
        src = cv2.imread(dst, cv2.IMREAD_COLOR)
        if src is not None:
            resized = src
            resize_ref = ""
            #리사이즈
            if src.shape[0] > image_resize_size[0]:
                if src.shape[1] > image_resize_size[1]: #높이도 기준보다 크고, 너비도 기준보다 큰 경우  → 높이 기준으로 리사이즈
                    resize_ref = "width"
                else: #높이는 기준보다 크지만, 너비는 기준보다 작은 경우 → 높이 기준으로 리사이즈
                    resize_ref = "height"
            else:
                if src.shape[1] > image_resize_size[1]: #높이는 기준보다 적고, 너비는 기준보다 큰 경우 → 너비 기준으로 리사이즈
                    resize_ref = "width"
                else: #높이가 기준보다 적고, 너비도 기준보다 작은 경우 → 리사이즈 불필요
                    pass

            if resize_ref is "width":
                resized=cv2.resize(src, dsize=(int(image_resize_size[1]), int(src.shape[0]*(image_resize_size[1]/src.shape[1]))), interpolation=cv2.INTER_AREA)
            elif resize_ref is "height":
                resized=cv2.resize(src, dsize=(int(src.shape[1]*(image_resize_size[0]/src.shape[0])), int(image_resize_size[0])), interpolation=cv2.INTER_AREA)

            #resized = faceRecognition(resized)

            #먼저 너비를 기준으로 조정 너비=4, 높이를 3에 맞춤
            width = resized.shape[1]
            height = int((width*4)/3)
            crop_axis = "height"

            #조정해야 할 높이가 현재 높이보다 크다면 높이를 기준으로 조정함
            if height > resized.shape[0]: 
                height = resized.shape[0]
                width = int((height*3)/4)
                crop_axis = "width"
            
            #잘라내야 할 높이 혹은 너비는?
            if crop_axis is "height":
                crop_size = resized.shape[0] - height
                resized = resized[0:height,:] #아래만 자른다
            elif crop_axis is "width":
                crop_size = int((resized.shape[1] - width)/2)
                crop_size_left = crop_size
                crop_size_right = resized.shape[1]-crop_size
                if (resized.shape[1] - width) % 2 > 0:
                    crop_size_left = crop_size + 1
                resized = resized[:,crop_size_left:crop_size_right] #좌우를 동등하게 자른다

            # print("최종 너비 : {}".format(width)+", 최종 높이 : {}".format(height)+", 결과 너비 : {}".format(resized.shape[1])+", 결과 높이 : {}".format(resized.shape[0])+", 잘라내야 할 축 : "+crop_axis+", 잘라내야 할 사이즈 : "+str(crop_size))

            cv2.imwrite(resized_path+"temp_resized"+"."+f.split(".")[1], resized)
            os.rename("./"+resized_path+"temp_resized"+"."+f.split(".")[1], "./"+resized_path+f.split(".")[0]+"_"+str(resized.shape[1])+"x"+str(resized.shape[0])+"."+f.split(".")[1])
        os.remove("./"+download_path+"temp")

def makePPT(resized_path):
    print("교육생 명단을 PPT로 작성하는 중...")
    df = pd.read_clipboard()
    if len(df) < 1:
        printInitError()
        sys.exit()

    prs = Presentation("master.pptx")
    slide = prs.slides[0]
    for student in df.iloc:
        # print(student["성명"]+" "+str(student["조"])+"-"+str(student["출력순서"]))
        for shape in slide.shapes:
            if shape.has_text_frame:
                this_paragraph = shape.text_frame.paragraphs[0]
                if str(student["조"]).strip() == this_paragraph.text.strip()[2:3] and str(student["출력순서"]).strip() == this_paragraph.text.strip()[4:5]:
                    this_label = this_paragraph.text.strip()[0:2]
                    if this_label == "사진":
                        files = os.listdir("./"+resized_path)
                        for f in files:
                            if student["휴대폰"] in f:
                                slide.shapes.add_picture("./"+resized_path+f, shape.left, shape.top, shape.width, shape.height)
                                break
                        this_paragraph.text = ""

            elif shape.has_table:
                for i in range(0, 7):
                    cell = shape.table.rows[i].cells[0]
                    this_table_paragraph = cell.text_frame.paragraphs[0]
                    this_label = this_table_paragraph.text.strip()[0:2]
                    if str(student["조"]).strip() == this_table_paragraph.text.strip()[2:3] and str(student["출력순서"]).strip() == this_table_paragraph.text.strip()[4:5]:
                        if this_label == "이름":
                            this_table_paragraph.text = student["성명"]
                            this_table_paragraph.font.bold = True
                        elif this_label == "나이":
                            this_table_paragraph.text = str(student["나이"])
                        elif this_label == "대학":
                            this_table_paragraph.text = student["대학명"]
                        elif this_label == "전공":
                            this_table_paragraph.text = student["학부전공"]
                        elif this_label == "지역":
                            this_table_paragraph.text = student["거주지"]
                        elif this_label == "전화":
                            this_table_paragraph.text = student["휴대폰"]
                        elif this_label == "숙소":
                            this_table_paragraph.text = student["숙소"]
                        this_table_paragraph.font.size = Pt(8)
                    
    prs.save("result.pptx")

if __name__ == "__main__":
    initCheck = checkInit()
    if initCheck:
        makeDownloadDirectory(image_path)
        makeDownloadDirectory(download_path)
        downloadStudentImages(login_id, login_pw, cha_name)
        cropImages(download_path)
        makePPT(resized_path)