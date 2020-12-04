import os
import six
import copy
import pandas as pd
from time import sleep
from pptx.util import Pt
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import ColorFormat, RGBColor

#폴더 생성 및 PPT 생성을 위한 정보
filesave_root_path = "results/"

def checkClipboard():
    print("===================================================================================================")     
    _ = input("데이터 복사가 완료되면 엔터키를 눌러 주세요.")
    try:
        df = pd.read_clipboard()
        if df is None or len(df) < 1 or "비즈니스 시나리오" not in df.columns:
            print("데이터가 클립보드로 복사되지 않았습니다.")
            sleep(2)
            checkClipboard()
    except:
        print("오류로 인해 데이터가 복사되지 않았습니다.")
        sleep(2)
        checkClipboard()
    return df

def makePPT():
    df = checkClipboard()
    print("PPT로 작성하는 중...")

    prs = Presentation("master.pptx")

    for idx in range(len(df)):
        duplicate_slide(prs, 1)
        this_slide = prs.slides[idx+2]
        for shape in this_slide.shapes:
            if shape.has_text_frame and shape.text_frame.text == "타이틀":
                shape.text_frame.paragraphs[0].text = str(df.loc[idx, "비즈니스 시나리오"])
                shape.text_frame.paragraphs[0].font.size = Pt(20)
                shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                shape.text_frame.paragraphs[0].font.bold = True
                shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            elif shape.has_table:
                for r in range(len(shape.table.rows)):
                    for c in range(len(shape.table.rows[r].cells)):
                        this_cell = shape.table.rows[r].cells[c]
                        this_table_paragraph = this_cell.text_frame.paragraphs[0]
                        for col_name in df:
                            if col_name == this_table_paragraph.text[:len(this_table_paragraph.text)-2] and this_table_paragraph.text[len(this_table_paragraph.text)-2:len(this_table_paragraph.text)] == "내용":
                                this_val = str(df.loc[idx, col_name])
                                if this_val == "" or this_val == "nan":
                                    this_val = "-"
                                if col_name == "정상구분" and this_val == "-":
                                    this_val = "정상"
                                this_table_paragraph.text = this_val.replace("*****br*****", " ")
                                this_table_paragraph.font.size = Pt(13)
                                this_table_paragraph.font.color.rgb = RGBColor(0, 0, 0)
                                this_table_paragraph.font.bold = False

    makeDownloadDirectory()
    prs.save(filesave_root_path+"비즈니스 시나리오.pptx")
    print("작성완료")

def duplicate_slide(prs, index):
    template = prs.slides[index]
    blank_slide_layout = prs.slide_layouts[0]
    copied_slide = prs.slides.add_slide(blank_slide_layout)

    for shp in template.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for _, value in six.iteritems(template.part.rels):
        # Make sure we don't copy a notesSlide relation as that won't exist
        if "notesSlide" not in value.reltype:
            copied_slide.part.rels.add_relationship(
                value.reltype,
                value._target,
                value.rId
            )

def makeDownloadDirectory():
    try:
        if not(os.path.isdir("./"+filesave_root_path)):
            os.makedirs(os.path.join("./"+filesave_root_path))
    except OSError as e:
        if e.errno != errno.EEXIST:
            print(filesave_root_path + " : 폴더 생성에 실패하였습니다.")
            raise

if __name__ == "__main__":
    makePPT()